#!/usr/bin/env python3
"""Generate Task2_Presentation.pptx for the IoT-23 cybersecurity team project."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

RESULTS = "results"
OUT = "Task2_Presentation.pptx"

# ---- Color palette (modern dark/cyber theme) -----------------------------
NAVY      = RGBColor(0x0A, 0x1F, 0x3D)   # deep navy (titles, dark backgrounds)
NAVY_DARK = RGBColor(0x05, 0x12, 0x26)   # near-black navy
TEAL      = RGBColor(0x00, 0xC8, 0xC8)   # bright cyan accent
RED       = RGBColor(0xE7, 0x4C, 0x3C)   # alert red
SOFT_BG   = RGBColor(0xF6, 0xF8, 0xFB)   # light slate background
SLATE     = RGBColor(0x47, 0x55, 0x69)   # secondary text
DARK_TEXT = RGBColor(0x1A, 0x20, 0x2C)   # primary text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---- Helpers --------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=DARK_TEXT, align=None,
             font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        if align:
            p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=16,
                color=DARK_TEXT, marker_color=TEAL):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Use a styled marker run + text run for color contrast
        run1 = p.add_run() if i > 0 else p.runs[0] if p.runs else None
        if run1 is None:
            run1 = p.add_run()
        run1.text = "▸  "
        run1.font.size = Pt(size)
        run1.font.bold = True
        run1.font.color.rgb = marker_color
        run1.font.name = "Calibri"

        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
        p.space_after = Pt(10)


def slide_chrome(slide, page_num, total, dark=False):
    """Add a footer / page number / project label to every slide."""
    bg = NAVY_DARK if dark else SOFT_BG
    fg = WHITE if dark else SLATE

    # full-bleed background
    add_rect(slide, 0, 0, SW, SH, bg)

    # top accent line
    add_rect(slide, 0, 0, SW, Inches(0.12), TEAL)

    # bottom footer bar
    add_rect(slide, 0, SH - Inches(0.45), SW, Inches(0.45),
             NAVY if not dark else NAVY)
    # Footer text
    add_text(slide, 0.4, 7.13, 9, 0.32,
             "IoT Botnet Detection  •  Linux Team Project  •  Alabama A&M University",
             size=10, color=WHITE)
    add_text(slide, 12.0, 7.13, 1.2, 0.32,
             f"{page_num:02d} / {total:02d}",
             size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)


def section_title(slide, kicker, title):
    """Standard content-slide header with kicker label + big title."""
    # Left vertical accent bar
    add_rect(slide, Inches(0.4), Inches(0.55), Inches(0.12), Inches(0.95), TEAL)

    add_text(slide, 0.7, 0.55, 12, 0.35, kicker.upper(),
             size=12, bold=True, color=TEAL)
    add_text(slide, 0.7, 0.85, 12.5, 0.75, title,
             size=30, bold=True, color=NAVY)
    # divider line
    add_rect(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.025), BORDER)


def stat_card(slide, left, top, width, height, value, label, accent=TEAL):
    add_rect(slide, left, top, width, height, WHITE, line=BORDER)
    # accent strip top
    add_rect(slide, left, top, width, Inches(0.08), accent)
    # value
    add_text(slide, left.inches + 0.1, top.inches + 0.25,
             width.inches - 0.2, 0.9,
             value, size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # label
    add_text(slide, left.inches + 0.1, top.inches + 1.05,
             width.inches - 0.2, 0.4,
             label.upper(), size=11, bold=True, color=SLATE, align=PP_ALIGN.CENTER)


# Total slide count for footer (we know it ahead)
TOTAL = 10

# ============================================================
# Slide 1 - Title (dark hero slide)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 1, TOTAL, dark=True)

# Big accent shape on the right
add_rect(s, Inches(9.5), Inches(0), Inches(3.83), Inches(7.5), NAVY)
# Vertical accent line
add_rect(s, Inches(9.4), Inches(0), Inches(0.08), Inches(7.5), TEAL)

# Kicker
add_text(s, 0.8, 1.6, 10, 0.4, "CYBERSECURITY  •  IoT-23 DATASET",
         size=14, bold=True, color=TEAL)

# Hero title
add_text(s, 0.8, 2.05, 11, 1.4, "Detecting IoT",
         size=64, bold=True, color=WHITE)
add_text(s, 0.8, 3.05, 11, 1.4, "Botnet Attacks",
         size=64, bold=True, color=TEAL)

# Subtitle
add_text(s, 0.8, 4.55, 11, 0.6,
         "A machine learning pipeline for classifying real-world malware traffic",
         size=18, color=RGBColor(0xB8, 0xC4, 0xD6))

# Team divider
add_rect(s, Inches(0.8), Inches(5.5), Inches(2.5), Inches(0.04), TEAL)
add_text(s, 0.8, 5.6, 12, 0.35, "TEAM", size=11, bold=True, color=TEAL)
add_text(s, 0.8, 5.95, 12, 0.4,
         "Thabo Ibrahim Traore  •  Jordan Barnes  •  Zizwe Mtonga  •  Tapiwa Musinga",
         size=15, color=WHITE)


# ============================================================
# Slide 2 - Objective
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 2, TOTAL)
section_title(s, "01  •  Objective", "What problem are we solving?")

# Two column layout - "The Problem" / "Our Goal"
# Left card
add_rect(s, Inches(0.5), Inches(1.95), Inches(6.2), Inches(4.7),
         WHITE, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.95), Inches(6.2), Inches(0.08), RED)
add_text(s, 0.75, 2.15, 5.8, 0.45, "THE PROBLEM",
         size=12, bold=True, color=RED)
add_text(s, 0.75, 2.55, 5.8, 0.6, "IoT devices are easy targets",
         size=20, bold=True, color=NAVY)
add_text(s, 0.75, 3.25, 5.8, 3.2,
         "Smart cameras, routers, and home gear get hijacked into "
         "botnets and used for DDoS attacks, data theft, and malware "
         "propagation. Detecting infections by hand is impossible at "
         "the scale of modern networks — automation is required.",
         size=14, color=SLATE)

# Right card
add_rect(s, Inches(6.85), Inches(1.95), Inches(6.0), Inches(4.7),
         WHITE, line=BORDER)
add_rect(s, Inches(6.85), Inches(1.95), Inches(6.0), Inches(0.08), TEAL)
add_text(s, 7.1, 2.15, 5.6, 0.45, "OUR GOAL",
         size=12, bold=True, color=TEAL)
add_text(s, 7.1, 2.55, 5.6, 0.6, "Automate intrusion detection",
         size=20, bold=True, color=NAVY)

goals = [
    "Classify network connections as benign or malicious",
    "Identify which features best predict attacks",
    "Run the pipeline on the ASAX supercomputer",
    "Hit high accuracy on real-world malware traffic",
]
add_bullets(s, 7.1, 3.25, 5.6, 3.2, goals, size=14)


# ============================================================
# Slide 3 - Dataset Overview
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 3, TOTAL)
section_title(s, "02  •  Dataset", "IoT-23 — Real Malware Traffic")

# Left text column
add_text(s, 0.7, 1.95, 6.5, 0.5,
         "23 IoT infection scenarios captured by Stratosphere Lab",
         size=16, bold=True, color=NAVY)

bullets_data = [
    "Source: Stratosphere Laboratory, Czech Technical University",
    "1.88 million labeled network connections",
    "Captured from real malware — not simulated",
    "Includes Mirai, Torii, and Hajime botnets",
    "23 capture scenarios + benign device baseline",
    "37 features per connection (flow stats, protocol, ports)",
]
add_bullets(s, 0.7, 2.65, 6.5, 4.2, bullets_data, size=14)

# Right column - stat cards (2x2 grid)
stat_card(s, Inches(7.6), Inches(1.95), Inches(2.6), Inches(1.55),
          "1.88M", "connections")
stat_card(s, Inches(10.4), Inches(1.95), Inches(2.6), Inches(1.55),
          "23", "scenarios")
stat_card(s, Inches(7.6), Inches(3.7), Inches(2.6), Inches(1.55),
          "37", "features")
stat_card(s, Inches(10.4), Inches(3.7), Inches(2.6), Inches(1.55),
          "192 MB", "data size")

# URL strip
add_rect(s, Inches(7.6), Inches(5.55), Inches(5.4), Inches(0.55), NAVY)
add_text(s, 7.7, 5.62, 5.3, 0.4,
         "stratosphereips.org/datasets-iot23",
         size=13, color=TEAL, bold=True)


# ============================================================
# Slide 4 - Attack Categories
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 4, TOTAL)
section_title(s, "03  •  Threat Landscape", "What we're detecting")

data = [
    ("Benign",   "120,499",   "Normal IoT device traffic"),
    ("DDoS",     "1,643,228", "Distributed denial-of-service flood attacks"),
    ("PortScan", "74,139",    "Scanning for open ports / vulnerable services"),
    ("C&C",      "22,502",    "Command & Control - infected device 'phones home'"),
    ("Botnet",   "15,252",    "Mirai / Okiru botnet propagation"),
    ("Attack",   "9,366",     "Other generic attack categories"),
]

# Custom-styled table card
add_rect(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(4.7),
         WHITE, line=BORDER)

# Header row
add_rect(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.55), NAVY)
add_text(s, 0.85, 2.02, 2.5, 0.4, "CATEGORY",
         size=12, bold=True, color=TEAL)
add_text(s, 4.0, 2.02, 3.0, 0.4, "CONNECTIONS",
         size=12, bold=True, color=TEAL)
add_text(s, 7.5, 2.02, 5.0, 0.4, "DESCRIPTION",
         size=12, bold=True, color=TEAL)

# Data rows
row_h = 0.65
for i, (cat, cnt, desc) in enumerate(data):
    y = 2.55 + i * row_h
    if i % 2 == 0:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3),
                 Inches(row_h), SOFT_BG)
    is_benign = (cat == "Benign")
    color = TEAL if is_benign else RED
    # Category badge
    add_rect(s, Inches(0.85), Inches(y + 0.13), Inches(0.18), Inches(0.4),
             color)
    add_text(s, 1.15, y + 0.15, 2.5, 0.4, cat,
             size=14, bold=True, color=NAVY)
    add_text(s, 4.0, y + 0.15, 3.0, 0.4, cnt,
             size=14, bold=True, color=SLATE)
    add_text(s, 7.5, y + 0.15, 5.2, 0.4, desc,
             size=13, color=SLATE)

# Caption
add_text(s, 0.5, 6.7, 12.3, 0.35,
         "⚠  HEAVILY IMBALANCED  —  93.6% of traffic is malicious (mostly DDoS)",
         size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 5 - Methodology
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 5, TOTAL)
section_title(s, "04  •  Methodology", "The data pipeline")

steps = [
    ("01", "Load",       "Read 23 CSV files, combine into 1.88M-row dataset"),
    ("02", "Preprocess", "Drop dead columns, fill missing values, encode labels"),
    ("03", "Explore",    "Generate distribution charts and protocol breakdowns"),
    ("04", "Split",      "70% train (1.32M) / 30% test (565K), stratified"),
    ("05", "Train",      "Random Forest classifier — 100 decision trees"),
    ("06", "Evaluate",   "Accuracy, precision, recall, confusion matrix"),
    ("07", "Visualize",  "Save all results as PNGs in results/"),
]

card_w = 1.75
gap = 0.05
start_x = 0.45
y_top = 2.05

for i, (num, label, desc) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    add_rect(s, Inches(x), Inches(y_top), Inches(card_w), Inches(3.0),
             WHITE, line=BORDER)
    # Top number band
    add_rect(s, Inches(x), Inches(y_top), Inches(card_w), Inches(0.55), NAVY)
    add_text(s, x + 0.1, y_top + 0.07, card_w - 0.2, 0.45,
             num, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    # Label
    add_text(s, x + 0.1, y_top + 0.7, card_w - 0.2, 0.45,
             label, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Description
    add_text(s, x + 0.1, y_top + 1.2, card_w - 0.2, 1.7,
             desc, size=10, color=SLATE, align=PP_ALIGN.CENTER)

# Tools strip
add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.7), NAVY)
add_text(s, 0.7, 5.5, 3.5, 0.5, "TECH STACK",
         size=11, bold=True, color=TEAL)
add_text(s, 4.0, 5.52, 9, 0.5,
         "Python 3   •   pandas   •   scikit-learn   •   matplotlib   •   seaborn",
         size=14, bold=True, color=WHITE)

# Note
add_text(s, 0.5, 6.4, 12.3, 0.35,
         "End-to-end runtime on full dataset: under 30 seconds",
         size=12, color=SLATE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 6 - Exploratory Data Analysis
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 6, TOTAL)
section_title(s, "05  •  Exploration", "What does the data look like?")

# Two image cards
img_top = 1.95
img_h = 4.7

add_rect(s, Inches(0.5), Inches(img_top), Inches(6.2), Inches(img_h),
         WHITE, line=BORDER)
add_rect(s, Inches(0.5), Inches(img_top), Inches(6.2), Inches(0.08), TEAL)
if os.path.exists(f"{RESULTS}/attack_distribution.png"):
    s.shapes.add_picture(f"{RESULTS}/attack_distribution.png",
                         Inches(0.65), Inches(img_top + 0.2),
                         width=Inches(5.9))

add_rect(s, Inches(6.85), Inches(img_top), Inches(6.0), Inches(img_h),
         WHITE, line=BORDER)
add_rect(s, Inches(6.85), Inches(img_top), Inches(6.0), Inches(0.08), TEAL)
if os.path.exists(f"{RESULTS}/benign_vs_malicious.png"):
    s.shapes.add_picture(f"{RESULTS}/benign_vs_malicious.png",
                         Inches(7.4), Inches(img_top + 0.2),
                         height=Inches(4.4))

# Caption
add_text(s, 0.5, 6.7, 12.3, 0.35,
         "Benign traffic is the minority — model must learn to spot the rare case",
         size=12, color=SLATE, align=PP_ALIGN.CENTER, bold=True)


# ============================================================
# Slide 7 - Results / Performance
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 7, TOTAL)
section_title(s, "06  •  Results", "Random Forest performance")

# Confusion matrix card on left
add_rect(s, Inches(0.5), Inches(1.95), Inches(6.5), Inches(4.95),
         WHITE, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.95), Inches(6.5), Inches(0.08), TEAL)
if os.path.exists(f"{RESULTS}/confusion_matrix.png"):
    s.shapes.add_picture(f"{RESULTS}/confusion_matrix.png",
                         Inches(0.85), Inches(2.15),
                         height=Inches(4.6))

# Right side - hero accuracy + metric grid
# Hero accuracy
add_rect(s, Inches(7.2), Inches(1.95), Inches(5.7), Inches(1.85), NAVY)
add_text(s, 7.3, 2.05, 5.5, 0.35, "OVERALL ACCURACY",
         size=11, bold=True, color=TEAL)
add_text(s, 7.3, 2.4, 5.5, 1.2, "97.87%",
         size=64, bold=True, color=WHITE)
add_text(s, 7.3, 3.45, 5.5, 0.35,
         "565,502 test connections classified correctly",
         size=11, color=RGBColor(0xB8, 0xC4, 0xD6))

# Metric grid (2x2)
mx, my = 7.2, 4.0
mw, mh = 2.75, 1.4
metrics_grid = [
    ("0.83", "BENIGN F1",       (0, 0)),
    ("0.99", "MALICIOUS F1",    (1, 0)),
    ("23.4s", "TRAINING TIME",   (0, 1)),
    ("0.7s",  "PREDICTION TIME", (1, 1)),
]
for val, label, (cx, cy) in metrics_grid:
    x = mx + cx * (mw + 0.2)
    y = my + cy * (mh + 0.15)
    add_rect(s, Inches(x), Inches(y), Inches(mw), Inches(mh),
             WHITE, line=BORDER)
    add_rect(s, Inches(x), Inches(y), Inches(mw), Inches(0.06), TEAL)
    add_text(s, x + 0.1, y + 0.18, mw - 0.2, 0.7,
             val, size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.95, mw - 0.2, 0.35,
             label, size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 8 - Feature Importances
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 8, TOTAL)
section_title(s, "07  •  Insights", "Which features matter most?")

# Left - chart card
add_rect(s, Inches(0.5), Inches(1.95), Inches(7.5), Inches(4.95),
         WHITE, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.95), Inches(7.5), Inches(0.08), TEAL)
if os.path.exists(f"{RESULTS}/feature_importances.png"):
    s.shapes.add_picture(f"{RESULTS}/feature_importances.png",
                         Inches(0.7), Inches(2.15),
                         height=Inches(4.6))

# Right - insight card
add_rect(s, Inches(8.2), Inches(1.95), Inches(4.7), Inches(4.95), NAVY)
add_text(s, 8.4, 2.1, 4.4, 0.4, "KEY INSIGHT",
         size=11, bold=True, color=TEAL)
add_text(s, 8.4, 2.5, 4.4, 0.65, "Volume tells the story",
         size=20, bold=True, color=WHITE)

# Highlight box
add_rect(s, Inches(8.4), Inches(3.3), Inches(4.3), Inches(1.5),
         RGBColor(0x12, 0x2B, 0x55))
add_text(s, 8.55, 3.4, 4.0, 0.4, "TOP PREDICTORS",
         size=10, bold=True, color=TEAL)
add_text(s, 8.55, 3.7, 4.0, 1.1,
         "orig_pkts\norig_ip_bytes\nresp_bytes",
         size=14, bold=True, color=WHITE, font="Consolas")

add_text(s, 8.4, 4.95, 4.4, 1.85,
         "DDoS and port-scan attacks send unusual packet volumes; "
         "C&C traffic shows distinctive byte signatures. The model "
         "rediscovered what real-world security analysts already know.",
         size=12, color=RGBColor(0xB8, 0xC4, 0xD6))


# ============================================================
# Slide 9 - ASAX Workflow
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 9, TOTAL)
section_title(s, "08  •  Deployment", "Running on the ASAX supercomputer")

# Left - bullet list of what the script does
add_text(s, 0.7, 1.95, 7.5, 0.45, "One shell script handles everything",
         size=18, bold=True, color=NAVY)

steps = [
    "Loads Python modules via 'module load python3'",
    "Installs dependencies (pandas, scikit-learn, matplotlib...)",
    "Auto-downloads IoT-23 dataset from Kaggle",
    "Combines 23 CSV files into a single dataset",
    "Runs the Python analysis pipeline",
    "Saves all charts and reports to results/ folder",
]
add_bullets(s, 0.7, 2.6, 7.5, 4.0, steps, size=14)

# Right - terminal-style code blocks
# Block 1: SLURM submission
term_x = 8.4
term_w = 4.5

add_rect(s, Inches(term_x), Inches(1.95), Inches(term_w), Inches(2.0),
         NAVY_DARK, line=BORDER)
# fake "window" header
add_rect(s, Inches(term_x), Inches(1.95), Inches(term_w), Inches(0.35),
         RGBColor(0x1E, 0x2B, 0x40))
add_text(s, term_x + 0.15, 1.97, 4.0, 0.3, "● ● ●  asax terminal",
         size=10, color=RGBColor(0x9C, 0xA8, 0xB8))

add_text(s, term_x + 0.2, 2.45, term_w - 0.3, 0.35,
         "# submit as SLURM job",
         size=11, color=RGBColor(0x6B, 0x82, 0xA0), font="Consolas")
add_text(s, term_x + 0.2, 2.78, term_w - 0.3, 0.4,
         "$ sbatch run_analysis.sh",
         size=14, bold=True, color=TEAL, font="Consolas")
add_text(s, term_x + 0.2, 3.2, term_w - 0.3, 0.35,
         "Submitted batch job 84291",
         size=11, color=RGBColor(0xB8, 0xC4, 0xD6), font="Consolas")

# Block 2: Interactive run
add_rect(s, Inches(term_x), Inches(4.15), Inches(term_w), Inches(2.0),
         NAVY_DARK, line=BORDER)
add_rect(s, Inches(term_x), Inches(4.15), Inches(term_w), Inches(0.35),
         RGBColor(0x1E, 0x2B, 0x40))
add_text(s, term_x + 0.15, 4.17, 4.0, 0.3, "● ● ●  local terminal",
         size=10, color=RGBColor(0x9C, 0xA8, 0xB8))

add_text(s, term_x + 0.2, 4.65, term_w - 0.3, 0.35,
         "# or run interactively",
         size=11, color=RGBColor(0x6B, 0x82, 0xA0), font="Consolas")
add_text(s, term_x + 0.2, 4.98, term_w - 0.3, 0.4,
         "$ bash run_analysis.sh",
         size=14, bold=True, color=TEAL, font="Consolas")
add_text(s, term_x + 0.2, 5.4, term_w - 0.3, 0.35,
         "Accuracy: 0.9787  ✓",
         size=11, color=RGBColor(0xB8, 0xC4, 0xD6), font="Consolas")


# ============================================================
# Slide 10 - Conclusion
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_chrome(s, 10, TOTAL, dark=True)

# Vertical accent bar
add_rect(s, Inches(0.4), Inches(0.55), Inches(0.12), Inches(0.95), TEAL)
add_text(s, 0.7, 0.55, 12, 0.35, "09  •  WRAP UP",
         size=12, bold=True, color=TEAL)
add_text(s, 0.7, 0.85, 12.5, 0.75, "Conclusion & Demo",
         size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.025),
         RGBColor(0x1E, 0x2B, 0x40))

# Three result cards horizontally
card_top = 2.05
card_h = 2.4
card_w = 4.0
gap_c = 0.25
start = 0.5

results_list = [
    ("97.87%", "Accuracy",   "On 565K real test connections from real malware"),
    ("1.88M",  "Connections", "Network flows analyzed end-to-end in under 30s"),
    ("100%",   "Automated",  "Single shell command runs the entire pipeline"),
]
for i, (val, lbl, desc) in enumerate(results_list):
    x = start + i * (card_w + gap_c)
    add_rect(s, Inches(x), Inches(card_top), Inches(card_w), Inches(card_h),
             NAVY)
    add_rect(s, Inches(x), Inches(card_top), Inches(card_w), Inches(0.08), TEAL)
    add_text(s, x + 0.2, card_top + 0.25, card_w - 0.4, 0.4, lbl.upper(),
             size=11, bold=True, color=TEAL)
    add_text(s, x + 0.2, card_top + 0.6, card_w - 0.4, 1.1, val,
             size=44, bold=True, color=WHITE)
    add_text(s, x + 0.2, card_top + 1.65, card_w - 0.4, 0.7, desc,
             size=12, color=RGBColor(0xB8, 0xC4, 0xD6))

# Why it matters
add_text(s, 0.5, 4.85, 12.3, 0.4, "WHY IT MATTERS",
         size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.25, 12.3, 0.7,
         "IoT security is one of the fastest growing threats in cybersecurity.",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.65, 12.3, 0.7,
         "The same techniques shown here protect real networks from real attacks.",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Demo callout
add_rect(s, Inches(2.5), Inches(6.45), Inches(8.3), Inches(0.55), TEAL)
add_text(s, 2.5, 6.52, 8.3, 0.4,
         "▸  LIVE DEMO    github.com/ibrahimt39/AAMU-Linux-Project",
         size=14, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
